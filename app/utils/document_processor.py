import io
import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
import xlrd
from pptx import Presentation
from bs4 import BeautifulSoup
import os
import tempfile
from pathlib import Path


class DocumentProcessor:
    """
    Multi-format document parser

    Supported formats:
    - PDF: PyPDF2 + pdfplumber (Vietnamese support)
    - DOCX/DOC: python-docx
    - XLSX: openpyxl
    - XLS: xlrd (Excel 97-2003)
    - PPTX: python-pptx
    - HTML: BeautifulSoup
    - TXT: Plain text

    Strategy:
    1. Detect file type từ extension/mime-type
    2. Parse theo format tương ứng
    3. Extract text + metadata
    4. Clean text (remove extra spaces, Vietnamese chars)
    5. Return normalized text
    """

    ALLOWED_EXTENSIONS = {
        ".pdf",
        ".docx",
        ".doc",
        ".xlsx",
        ".xls",
        ".pptx",
        ".html",
        ".txt",
    }
    ALLOWED_MIME_TYPES = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
        "text/html",
        "text/plain",
    }

    @staticmethod
    def validate_file(filename: str, mime_type: str) -> tuple[bool, str]:
        """
        Validate file format

        Returns: (is_valid, error_message)
        """
        # Check extension
        ext = Path(filename).suffix.lower()
        if ext not in DocumentProcessor.ALLOWED_EXTENSIONS:
            return (
                False,
                f"Unsupported file type: {ext}. Allowed: {', '.join(DocumentProcessor.ALLOWED_EXTENSIONS)}",
            )

        # Check MIME type
        if mime_type not in DocumentProcessor.ALLOWED_MIME_TYPES:
            return False, f"Invalid MIME type: {mime_type}"

        return True, ""

    @staticmethod
    def extract_text_from_pdf(file_content: bytes, filename: str) -> str:
        """
        Extract text từ PDF

        Strategy:
        1. Try pdfplumber (better Vietnamese support)
        2. Fallback to PyPDF2
        3. Preserve text structure (paragraphs, lists)

        Args:
            file_content: Raw PDF bytes
            filename: Filename (for logging)

        Returns:
            Extracted text
        """
        text = ""

        try:
            # Write to temp file (pdfplumber cần file path)
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name

            try:
                # Try pdfplumber (best for Vietnamese)
                with pdfplumber.open(tmp_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n--- Page {page_num + 1} ---\n"
                            text += page_text
            except Exception as e:
                print(f"pdfplumber failed: {e}, trying PyPDF2...")

                # Fallback to PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num + 1} ---\n"
                        text += page_text

            finally:
                # Clean temp file
                os.unlink(tmp_path)

        except Exception as e:
            raise Exception(f"Failed to extract PDF: {str(e)}")

        return text if text else ""

    @staticmethod
    def extract_text_from_docx(file_content: bytes) -> str:
        """
        Extract text từ DOCX/DOC

        Preserves:
        - Paragraphs
        - Headings
        - Lists
        - Tables
        - Vietnamese characters

        Args:
            file_content: Raw DOCX bytes

        Returns:
            Extracted text
        """
        text = ""

        try:
            import io

            docx_file = DocxDocument(io.BytesIO(file_content))

            # Extract paragraphs
            for para in docx_file.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

            # Extract tables
            for table in docx_file.tables:
                text += "\n--- TABLE ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"

        except Exception as e:
            raise Exception(f"Failed to extract DOCX: {str(e)}")

        return text if text else "Unable to extract text from DOCX"

    @staticmethod
    def extract_text_from_docx_2(file_content: bytes) -> str:
        import re, collections

        docx_file = DocxDocument(io.BytesIO(file_content))

        def is_bold(para):
            return any(run.bold for run in para.runs if run.text.strip())

        def detect_heading_with_text(para) -> tuple[str | None, str]:
            """Trả về (prefix, text_để_dùng_làm_heading)"""
            text = para.text.strip()
            if not text:
                return None, text

            bold = is_bold(para)
            is_numbered = bool(re.match(r"^([IVX]+\.|\d+\.|\d+\.\d+)\s+\S", text))

            # 1. BỘ LỌC NHANH: Không in đậm VÀ không đánh số -> Chắc chắn là text thường
            if not bold and not is_numbered:
                return None, text

            heading_text = text

            # 2. LUỒNG ƯU TIÊN: Xử lý đoạn có Đánh Số (bất kể có in đậm hay không)
            if is_numbered:
                # Cơ chế chống nhận diện nhầm đoạn văn dài thành Heading
                if len(text) > 150:
                    colon_pos = text.find(":")
                    if 0 < colon_pos < 120:
                        heading_text = text[:colon_pos].strip()
                    else:
                        return None, text
                
                # Chốt chặn phụ: Nếu đã cắt rồi mà vẫn quá dài, hoặc có dấu ':' sát đầu (như "1. Lưu ý:")
                if len(heading_text) > 150 or ":" in heading_text[:40]:
                    return None, text

                # Phân loại Heading theo cấp độ số
                if re.match(r"^[IVX]+\.\s+\S", heading_text):
                    return "#", heading_text
                if re.match(r"^\d+\.\s+\S", heading_text) and not re.match(r"^\d+\.\d+", heading_text):
                    return "##", heading_text
                if re.match(r"^\d+\.\d+[\d.]*\s+\S", heading_text):
                    return "###", heading_text

            # 3. LUỒNG BỔ SUNG: Không đánh số NHƯNG in đậm (Tiêu đề tự do)
            # Điều kiện: Phải ngắn (<= 100 ký tự) và không có dấu hai chấm ở đầu để tránh nhầm với câu nhấn mạnh
            if bold and len(text) <= 100 and ":" not in text[:40]:
                return "###", text # Mặc định gom tiêu đề dạng này thành H3 an toàn

            # 4. Rớt đài: Trả về text thường
            return None, text

        def table_to_markdown(table) -> str:
            rows = []
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            return "\n".join(rows)
        def table_to_semantic_text(table) -> str:
            if not table.rows:
                return ""
                
            # Lấy tiêu đề cột từ dòng đầu tiên
            headers = [cell.text.strip().replace("\n", " ") for cell in table.rows[0].cells]
            
            semantic_rows = []
            for i, row in enumerate(table.rows):
                if i == 0:
                    continue # Bỏ qua dòng header
                    
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                
                row_text_parts = []
                for j, cell_val in enumerate(cells):
                    # Nếu cột không có tiêu đề, gán tên tạm là Cột 1, Cột 2...
                    col_name = headers[j] if j < len(headers) and headers[j] else f"Cột {j+1}"
                    if cell_val: # Chỉ lấy các ô có dữ liệu để tiết kiệm token
                        row_text_parts.append(f"{col_name}: {cell_val}")
                
                if row_text_parts:
                    semantic_rows.append("- " + ", ".join(row_text_parts))
                    
            return "\n".join(semantic_rows)
        lines = []
        for child in docx_file.element.body.iterchildren():
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                from docx.text.paragraph import Paragraph

                para = Paragraph(child, docx_file)
                if not para.text.strip():
                    continue

                prefix, heading_text = detect_heading_with_text(para)
                if prefix:
                    lines.append(f"{prefix} {heading_text}")
                    # Nếu heading bị cắt, vẫn giữ phần còn lại làm body text
                    remaining = (
                        para.text.strip()[len(heading_text) :].strip().lstrip(":")
                    )
                    if remaining:
                        lines.append(remaining.strip())
                else:
                    lines.append(para.text.strip())

            elif tag == "tbl":
                from docx.table import Table

                #lines.append(table_to_markdown(Table(child, docx_file)))
                lines.append(table_to_semantic_text(Table(child, docx_file)))
                lines.append("")

        return "\n".join(lines)

    @staticmethod
    def extract_text_from_pdf_2(file_content: bytes, filename: str) -> str:
        import collections

        def _extract_with_pdfplumber(tmp_path: str) -> str:
            with pdfplumber.open(tmp_path) as pdf:
                # Bước 1: Thu thập font size toàn bộ document để tìm body_size
                size_counter = collections.Counter()
                for page in pdf.pages:
                    for w in page.extract_words(
                        extra_attrs=["size"], keep_blank_chars=False
                    ):
                        size = round(w.get("size", 0))
                        if size > 0:
                            size_counter[size] += 1

                if not size_counter:
                    return ""

                body_size = size_counter.most_common(1)[0][0]

                def to_heading(size: float) -> str | None:
                    diff = size - body_size
                    if diff >= 6:
                        return "#"
                    if diff >= 3:
                        return "##"
                    if diff >= 1.5:
                        return "###"
                    return None

                def table_to_markdown(table_data: list) -> str:
                    rows = []
                    for i, row in enumerate(table_data):
                        cells = [
                            (cell or "").replace("\n", " ").strip() for cell in row
                        ]
                        rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    return "\n".join(rows)
                def table_to_semantic_text(table_data: list) -> str:
                    if not table_data:
                        return ""
                        
                    # Dòng đầu tiên làm header
                    headers = [(cell or "").replace("\n", " ").strip() for cell in table_data[0]]
                    
                    semantic_rows = []
                    for i, row in enumerate(table_data):
                        if i == 0:
                            continue
                            
                        cells = [(cell or "").replace("\n", " ").strip() for cell in row]
                        
                        row_text_parts = []
                        for j, cell_val in enumerate(cells):
                            col_name = headers[j] if j < len(headers) and headers[j] else f"Cột {j+1}"
                            if cell_val:
                                row_text_parts.append(f"{col_name}: {cell_val}")
                                
                        if row_text_parts:
                            semantic_rows.append("- " + ", ".join(row_text_parts))
                            
                    return "\n".join(semantic_rows)

                all_output = []

                for page in pdf.pages:
                    # Lấy bounding box các bảng trong trang
                    tables = page.find_tables()
                    table_bboxes = [t.bbox for t in tables]

                    def in_table(word) -> bool:
                        x0, top, x1, bottom = (
                            word["x0"],
                            word["top"],
                            word["x1"],
                            word["bottom"],
                        )
                        for tx0, ttop, tx1, tbottom in table_bboxes:
                            if (
                                x0 >= tx0
                                and x1 <= tx1
                                and top >= ttop
                                and bottom <= tbottom
                            ):
                                return True
                        return False

                    # Extract words nằm ngoài bảng (heading + body text)
                    words = page.extract_words(
                        extra_attrs=["size"], keep_blank_chars=False
                    )
                    non_table_words = [w for w in words if not in_table(w)]

                    line_groups: dict[int, list] = {}
                    for w in non_table_words:
                        y_key = round(w["top"] / 3) * 3
                        line_groups.setdefault(y_key, []).append(w)

                    prev_heading_text = None
                    prev_heading_level = None
                    prev_heading_size = None

                    for y_key in sorted(line_groups):
                        row = sorted(line_groups[y_key], key=lambda w: w["x0"])
                        text = " ".join(w["text"] for w in row).strip()
                        if not text:
                            continue

                        first_size = round(row[0].get("size", body_size))
                        h = to_heading(first_size)

                        if (
                            h
                            and prev_heading_level == h
                            and prev_heading_size == first_size
                        ):
                            prev_heading_text = prev_heading_text + " " + text
                            all_output[-1] = f"{h} {prev_heading_text}"
                        elif h:
                            all_output.append(f"{h} {text}")
                            prev_heading_text = text
                            prev_heading_level = h
                            prev_heading_size = first_size
                        else:
                            all_output.append(text)
                            prev_heading_text = None
                            prev_heading_level = None
                            prev_heading_size = None

                    # Extract bảng theo cell (giữ nguyên nội dung, không bị wrap dòng)
                    for table in tables:
                        data = table.extract()
                        if not data:
                            continue
                        #all_output.append(table_to_markdown(data))
                        all_output.append(table_to_semantic_text(data))
                        all_output.append("")

                return "\n".join(all_output)

        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            try:
                text = _extract_with_pdfplumber(tmp_path)
                if text.strip():
                    return DocumentProcessor.clean_text(text)
            finally:
                os.unlink(tmp_path)
        except Exception as e:
            print(f"pdfplumber failed: {e}, fallback PyPDF2")

        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        lines = [p.extract_text() for p in pdf_reader.pages if p.extract_text()]
        return DocumentProcessor.clean_text("\n".join(lines))

    @staticmethod
    def extract_text_from_xlsx(file_content: bytes) -> str:
        """
        Extract text từ Excel

        Includes:
        - All sheets
        - Cell values
        - Preserves structure

        Args:
            file_content: Raw XLSX bytes

        Returns:
            Extracted text
        """
        text = ""

        try:
            import io

            workbook = load_workbook(io.BytesIO(file_content))

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== SHEET: {sheet_name} ===\n"

                for row in sheet.iter_rows():
                    row_text = " | ".join(
                        [
                            str(cell.value) if cell.value is not None else ""
                            for cell in row
                        ]
                    )
                    if row_text.strip():
                        text += row_text + "\n"

        except Exception as e:
            error_msg = str(e)
            if "File is not a zip file" in error_msg:
                raise Exception(
                    "Failed to extract XLSX: File is not a valid XLSX file. "
                    "This may be an old Excel 97-2003 (.xls) file renamed to .xlsx, "
                    "or the file is corrupted. Please use the correct extension."
                )
            raise Exception(f"Failed to extract XLSX: {error_msg}")

        return text if text else "Unable to extract text from XLSX"

    @staticmethod
    def extract_text_from_xls(file_content: bytes) -> str:
        """
        Extract text từ Excel 97-2003 (.xls) sử dụng xlrd

        Includes:
        - All sheets
        - Cell values
        - Preserves structure

        Args:
            file_content: Raw XLS bytes

        Returns:
            Extracted text
        """
        text = ""

        try:
            workbook = xlrd.open_workbook(file_contents=file_content)

            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text += f"\n=== SHEET: {sheet.name} ===\n"

                for row_idx in range(sheet.nrows):
                    row_values = [
                        str(cell) if cell is not None else ""
                        for cell in sheet.row_values(row_idx)
                    ]
                    row_text = " | ".join(row_values)
                    if row_text.strip():
                        text += row_text + "\n"

        except Exception as e:
            raise Exception(f"Failed to extract XLS: {str(e)}")

        return text if text else "Unable to extract text from XLS"

    @staticmethod
    def extract_text_from_pptx(file_content: bytes) -> str:
        """
        Extract text từ PowerPoint

        Includes:
        - Slide titles
        - Text boxes
        - Notes

        Args:
            file_content: Raw PPTX bytes

        Returns:
            Extracted text
        """
        text = ""

        try:
            import io

            prs = Presentation(io.BytesIO(file_content))

            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- SLIDE {slide_num + 1} ---\n"

                # Extract from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

                    # Extract from tables in slide
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            if row_text.strip():
                                text += row_text + "\n"

        except Exception as e:
            raise Exception(f"Failed to extract PPTX: {str(e)}")

        return text if text else "Unable to extract text from PPTX"

    @staticmethod
    def extract_text_from_html(file_content: bytes) -> str:
        """
        Extract text từ HTML

        Removes:
        - Scripts, styles
        - HTML tags
        - Extra whitespace

        Args:
            file_content: Raw HTML bytes

        Returns:
            Extracted text
        """
        text = ""

        try:
            html_str = file_content.decode("utf-8", errors="ignore")
            soup = BeautifulSoup(html_str, "html.parser")

            # Remove script and style
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text()

            # Clean up: remove multiple spaces/newlines
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

        except Exception as e:
            raise Exception(f"Failed to extract HTML: {str(e)}")

        return text if text else "Unable to extract text from HTML"

    @staticmethod
    def extract_text(file_content: bytes, filename: str, mime_type: str) -> str:
        """
        Main extraction function - route to specific handler

        Args:
            file_content: Raw file bytes
            filename: Filename
            mime_type: MIME type

        Returns:
            Extracted text (cleaned)
        """

        # Validate
        is_valid, error_msg = DocumentProcessor.validate_file(filename, mime_type)
        if not is_valid:
            raise ValueError(error_msg)

        ext = Path(filename).suffix.lower()

        # Route to appropriate handler
        if ext == ".pdf":
            text = DocumentProcessor.extract_text_from_pdf(file_content, filename)
        elif ext in [".docx", ".doc"]:
            text = DocumentProcessor.extract_text_from_docx(file_content)
        elif ext == ".xlsx":
            text = DocumentProcessor.extract_text_from_xlsx(file_content)
        elif ext == ".xls":
            text = DocumentProcessor.extract_text_from_xls(file_content)
        elif ext == ".pptx":
            text = DocumentProcessor.extract_text_from_pptx(file_content)
        elif ext == ".html":
            text = DocumentProcessor.extract_text_from_html(file_content)
        elif ext == ".txt":
            text = file_content.decode("utf-8", errors="ignore")
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        # Clean text
        text = DocumentProcessor.clean_text(text)

        return text

    @staticmethod
    def extract_text_ocr(file_content: bytes, filename: str) -> str:
        """
        OCR extraction for scanned PDFs using pymupdf + pytesseract.

        Args:
            file_content: Raw file bytes
            filename: Filename

        Returns:
            OCR'd text
        """
        import fitz
        import pytesseract
        from PIL import Image
        from app.core.config import settings

        pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD_PATH

        text_parts = []
        doc = fitz.open(stream=file_content, filetype="pdf")

        for idx, page in enumerate(doc):
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            page_text = pytesseract.image_to_string(img, lang="eng+vie")
            if page_text.strip():
                text_parts.append(f"\n--- Page {idx + 1} ---\n{page_text}")

        doc.close()
        result = "\n\n".join(text_parts)
        return DocumentProcessor.clean_text(result) if result else ""

    @staticmethod
    def clean_text(text: str) -> str:
        """
        Clean extracted text

        Remove:
        - Extra whitespace
        - Weird characters
        - Preserve Vietnamese accents

        Args:
            text: Raw text

        Returns:
            Cleaned text
        """
        import re

        # Remove multiple spaces
        text = re.sub(r" +", " ", text)

        # Xóa các dòng chỉ chứa dấu gạch ngang (---) sinh ra từ bảng biểu hoặc file Markdown
        text = re.sub(r"^-{3,}\s*$", "", text, flags=re.MULTILINE)
        # Remove multiple newlines
        text = re.sub(r"\n\n+", "\n\n", text)

        # Trim each line
        lines = [line.strip() for line in text.split("\n")]
        text = "\n".join(lines)

        # Remove leading/trailing whitespace
        text = text.strip()

        return text


documentProcessor = DocumentProcessor()
